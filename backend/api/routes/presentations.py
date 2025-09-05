from flask import request, jsonify, Blueprint, current_app, g, send_file, url_for
import io
import os
import requests
from pptx import Presentation as PptxPresentation
from .decorators import token_required
from pptx.util import Inches, Pt
import uuid
from PIL import Image
import ffmpeg
import win32com.client
import pythoncom
from ..models import Presentation, Slide
from ..extensions import db
from pptx.dml.color import RGBColor

presentations_bp = Blueprint('presentations', __name__)

PIXELS_PER_INCH = 80.0
SERVER_BASE_URL = 'http://127.0.0.1:5000'

def px_to_inches(px):
    return px / PIXELS_PER_INCH

def _download_image_from_url(image_url):
    if image_url.startswith('/'):
        image_url = f"{SERVER_BASE_URL}{image_url}"
    try:
        response = requests.get(image_url, stream=True, timeout=30)
        response.raise_for_status()
        image_data = io.BytesIO(response.content)
        image_data.seek(0)
        return image_data
    except Exception as e:
        print(f"Ошибка загрузки изображения {image_url}: {e}")
        return None

def _create_pptx_from_data(presentation_id):
    presentation_data = Presentation.query.get_or_404(presentation_id)
    
    prs = PptxPresentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    slides = Slide.query.filter_by(presentation_id=presentation_data.id).order_by(Slide.slide_number).all()

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if slide_data.background_image:
            try:
                print(f"Слайд {slide_data.slide_number}: пытаюсь применить картинку {slide_data.background_image}")
                
                image_stream = _download_image_from_url(slide_data.background_image)
                if image_stream:
                    pic = slide.shapes.add_picture(
                        image_stream, 
                        Inches(0), Inches(0), 
                        width=prs.slide_width, 
                        height=prs.slide_height
                    )
                    
                    slide.shapes._spTree.insert(2, pic._element)

            except Exception as e:
                print(f"Не удалось добавить фон-картинку {slide_data.background_image}: {e}")

        elif slide_data.background_color:
            try:
                print(f"Слайд {slide_data.slide_number}: пытаюсь применить цвет {slide_data.background_color}")

                hex_color = slide_data.background_color.lstrip('#')
                if len(hex_color) == 6:
                    r = int(hex_color[0:2], 16)
                    g = int(hex_color[2:4], 16)
                    b = int(hex_color[4:6], 16)
                    
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
            except Exception as e:
                print(f"Не удалось установить цвет фона {slide_data.background_color}: {e}")

        for element in slide_data.elements:
            container_left = Inches(px_to_inches(element.pos_x))
            container_top = Inches(px_to_inches(element.pos_y))
            container_width = Inches(px_to_inches(element.width))
            container_height = Inches(px_to_inches(element.height))

            if element.element_type == 'TEXT':
                txBox = slide.shapes.add_textbox(container_left, container_top, container_width, container_height)
                tf = txBox.text_frame
                tf.text = element.content or ""
                tf.word_wrap = True
                if tf.paragraphs:
                    tf.paragraphs[0].font.size = Pt(element.font_size or 24)
            
            elif element.element_type == 'IMAGE' and element.content:
                try:
                    image_stream = _download_image_from_url(element.content)
                    if not image_stream:
                        continue
                    
                    with Image.open(image_stream) as img:
                        img_width, img_height = img.size
                    
                    image_stream.seek(0)
                    
                    container_aspect = container_width.emu / container_height.emu if container_height.emu > 0 else 1
                    img_aspect = img_width / img_height if img_height > 0 else 1
                    
                    if img_aspect > container_aspect:
                        new_width = container_width
                        new_height = new_width / img_aspect
                    else:
                        new_height = container_height
                        new_width = new_height * img_aspect
                    
                    left_offset = (container_width - new_width) / 2
                    top_offset = (container_height - new_height) / 2
                    final_left = container_left + left_offset
                    final_top = container_top + top_offset
                    
                    slide.shapes.add_picture(image_stream, final_left, final_top, width=new_width, height=new_height)
                    
                except Exception as e:
                    print(f"Не удалось добавить изображение {element.content}: {e}")

            elif element.element_type == 'YOUTUBE_VIDEO' and element.content:
                image_stream = None
                thumbnail_urls = [
                    f"https://img.youtube.com/vi/{element.content}/maxresdefault.jpg",
                    f"https://img.youtube.com/vi/{element.content}/hqdefault.jpg", 
                    f"https://img.youtube.com/vi/{element.content}/0.jpg"
                ]
                for url in thumbnail_urls:
                    image_stream = _download_image_from_url(url)
                    if image_stream:
                        break
                
                if image_stream:
                    try:
                        pic = slide.shapes.add_picture(image_stream, container_left, container_top, width=container_width, height=container_height)
                        hlink = pic.click_action.hyperlink
                        hlink.address = f"https://www.youtube.com/watch?v={element.content}"
                    except Exception as e:
                        print(f"Не удалось добавить эскиз видео для {element.content}: {e}")
            
            elif (element.element_type == 'UPLOADED_VIDEO' or element.element_type == 'AUDIO') and element.content:
                try:
                    filename = element.content.split('/')[-1]
                    file_path = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
                    
                    if element.element_type == 'UPLOADED_VIDEO':
                        poster_path = os.path.join(current_app.root_path, 'static', 'video_poster.png')
                        mime_type = 'video/mp4'
                    else:
                        poster_path = os.path.join(current_app.root_path, 'static', 'audio_poster.png')
                        mime_type = 'audio/mpeg'

                    if os.path.exists(file_path) and os.path.exists(poster_path):
                        slide.shapes.add_movie(
                            file_path, container_left, container_top, 
                            container_width, container_height, 
                            poster_frame_image=poster_path, mime_type=mime_type
                        )
                    else:
                        if not os.path.exists(file_path):
                            print(f"Файл не найден для экспорта: {file_path}")
                        if not os.path.exists(poster_path):
                            print(f"Файл-заглушка (poster) не найден: {poster_path}")
                except Exception as e:
                    print(f"Не удалось добавить медиафайл {element.content}: {e}")
    
    return prs, presentation_data.title

@presentations_bp.route('/presentations/<string:presentation_id>/download/pptx', methods=['GET'])
@token_required
def download_presentation_pptx(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id: return jsonify({'message': 'Доступ запрещен'}), 403
    
    prs, title = _create_pptx_from_data(presentation_id)
    
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return send_file(file_stream, as_attachment=True, download_name=f"{title}.pptx", mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

@presentations_bp.route('/presentations/<string:presentation_id>/download/pdf', methods=['GET'])
@token_required
def download_presentation_pdf(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id:
        return jsonify({'message': 'Доступ запрещен'}), 403

    prs, title = _create_pptx_from_data(presentation_id)
    
    instance_path = current_app.instance_path
    os.makedirs(instance_path, exist_ok=True)
    
    pptx_filename = f"{uuid.uuid4()}.pptx"
    pdf_filename = f"{uuid.uuid4()}.pdf"
    
    pptx_path = os.path.join(instance_path, pptx_filename)
    pdf_path = os.path.join(instance_path, pdf_filename)
    
    powerpoint = None
    pres = None
    pdf_stream = None

    try:
        prs.save(pptx_path)

        pythoncom.CoInitializeEx(0)
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        pres = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        
        pres.SaveAs(pdf_path, 32)
        print(f"Successfully converted {pptx_path} to {pdf_path}")

        with open(pdf_path, 'rb') as f:
            pdf_stream = io.BytesIO(f.read())
        
        pdf_stream.seek(0)
        return send_file(pdf_stream, as_attachment=True, download_name=f"{title}.pdf", mimetype='application/pdf')

    except Exception as e:
        print(f"Failed to convert to PDF: {e}")
        return jsonify({"message": "Ошибка при конвертации в PDF. Убедитесь, что MS PowerPoint установлен на сервере."}), 500
    
    finally:
        if pres: pres.Close()
        if powerpoint: powerpoint.Quit()
        if os.path.exists(pptx_path): os.remove(pptx_path)
        if os.path.exists(pdf_path): os.remove(pdf_path)

def _serialize_slide(slide):
    if not slide:
        return None
    
    elements_output = []
    for e in slide.elements:
        element_data = {
            'id': e.id, 'element_type': e.element_type, 'pos_x': e.pos_x,
            'pos_y': e.pos_y, 'width': e.width, 'height': e.height,
            'content': e.content, 'font_size': e.font_size
        }
        if e.element_type == 'YOUTUBE_VIDEO':
            element_data['thumbnailUrl'] = f"https://img.youtube.com/vi/{e.content}/0.jpg"
        elements_output.append(element_data)
        
    return {
        'id': slide.id, 
        'slide_number': slide.slide_number,
        'background_color': slide.background_color, 
        'background_image': slide.background_image,
        'elements': elements_output
    }

@presentations_bp.route('/presentations', methods=['POST'])
@token_required
def create_presentation():
    data = request.get_json()
    title = data.get('title', 'Новая презентация')
    new_presentation = Presentation(title=title, owner=g.current_user)
    db.session.add(new_presentation)
    db.session.flush()

    first_slide = Slide(slide_number=1, presentation_id=new_presentation.id)
    db.session.add(first_slide)
    db.session.commit()
    
    first_slide_data = _serialize_slide(first_slide)
    
    return jsonify({
        'id': new_presentation.id,
        'title': new_presentation.title,
        'updated_at': new_presentation.updated_at.isoformat(),
        'first_slide': first_slide_data
    }), 201

@presentations_bp.route('/presentations/<string:presentation_id>', methods=['GET'])
@token_required
def get_presentation_by_id(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id:
        return jsonify({'message': 'Доступ запрещен'}), 403

    slides = Slide.query.filter_by(presentation_id=presentation.id).order_by(Slide.slide_number).all()
    slides_output = [_serialize_slide(s) for s in slides]

    return jsonify({'id': presentation.id, 'title': presentation.title, 'slides': slides_output}), 200

@presentations_bp.route('/presentations', methods=['GET'])
@token_required
def get_presentations():
    presentations = Presentation.query.filter_by(user_id=g.current_user.id, is_template=False).order_by(Presentation.updated_at.desc()).all()
    output = []
    for p in presentations:
        first_slide = Slide.query.filter_by(presentation_id=p.id, slide_number=1).first()
        first_slide_data = _serialize_slide(first_slide)
        
        output.append({
            'id': p.id,
            'title': p.title,
            'updated_at': p.updated_at.isoformat(),
            'first_slide': first_slide_data
        })
    return jsonify(output), 200

@presentations_bp.route('/presentations/<string:presentation_id>', methods=['DELETE'])
@token_required
def delete_presentation(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id: return jsonify({'message': 'Доступ запрещен'}), 403
    db.session.delete(presentation)
    db.session.commit()
    return jsonify({'message': 'Презентация успешно удалена'}), 200

@presentations_bp.route('/presentations/<string:presentation_id>', methods=['PUT'])
@token_required
def update_presentation(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != g.current_user.id: return jsonify({'message': 'Доступ запрещен'}), 403
    data = request.get_json()
    if 'title' in data: presentation.title = data['title']
    db.session.commit()
    
    first_slide = Slide.query.filter_by(presentation_id=presentation.id, slide_number=1).first()
    first_slide_data = _serialize_slide(first_slide)

    return jsonify({
        'id': presentation.id,
        'title': presentation.title,
        'updated_at': presentation.updated_at.isoformat(),
        'first_slide': first_slide_data
    }), 200

@presentations_bp.route('/upload/image', methods=['POST'])
@token_required
def upload_image():
    if 'file' not in request.files:
        return jsonify({'message': 'Файл не найден'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'message': 'Файл не выбран'}), 400
    if file:
        _root, extension = os.path.splitext(file.filename)
        filename = f"{uuid.uuid4()}{extension}"
        os.makedirs(current_app.config['UPLOAD_FOLDER'], exist_ok=True)
        save_path = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
        file.save(save_path)
        file_url = url_for('static', filename=f'uploads/{filename}', _external=False)
        return jsonify({'url': file_url}), 200

@presentations_bp.route('/upload/video', methods=['POST'])
@token_required
def upload_video():
    if 'file' not in request.files:
        return jsonify({'message': 'Файл не найден'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'message': 'Файл не выбран'}), 400
    if file:
        _root, extension = os.path.splitext(file.filename)
        temp_filename = f"{uuid.uuid4()}{extension}"
        final_filename = f"{uuid.uuid4()}.mp4"
        temp_path = os.path.join(current_app.config['UPLOAD_FOLDER'], temp_filename)
        final_path = os.path.join(current_app.config['UPLOAD_FOLDER'], final_filename)
        os.makedirs(current_app.config['UPLOAD_FOLDER'], exist_ok=True)
        file.save(temp_path)
        try:
            print(f"Starting conversion from {temp_path} to {final_path}")
            ffmpeg.input(temp_path).output(final_path, vcodec='libx24', acodec='aac', strict='experimental').run(capture_stdout=True, capture_stderr=True)
            print("Conversion successful")
            os.remove(temp_path)
            file_url = url_for('static', filename=f'uploads/{final_filename}', _external=False)
            return jsonify({'url': file_url}), 200
        except ffmpeg.Error as e:
            os.remove(temp_path)
            print("FFmpeg Error:")
            print(e.stderr.decode())
            return jsonify({'message': 'Не удалось обработать видеофайл'}), 500

@presentations_bp.route('/upload/audio', methods=['POST'])
@token_required
def upload_audio():
    if 'file' not in request.files:
        return jsonify({'message': 'Файл не найден'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'message': 'Файл не выбран'}), 400
    if file:
        _root, extension = os.path.splitext(file.filename)
        filename = f"{uuid.uuid4()}{extension}"
        os.makedirs(current_app.config['UPLOAD_FOLDER'], exist_ok=True)
        save_path = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
        file.save(save_path)
        file_url = url_for('static', filename=f'uploads/{filename}', _external=False)
        return jsonify({'url': file_url}), 200